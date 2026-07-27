from pptx import Presentation
from pptx.util import Pt

def create_presentation():
    prs = Presentation()

    # Slide 1: Title
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    title = slide.shapes.title
    subtitle = slide.placeholders[1]
    title.text = "FinSight N100 | Financial Intelligence Platform"
    subtitle.text = ("Enterprise-Grade Financial Analytics, REST API & BI Solution\n\n"
                     "10 Years Data | 92 Nifty 100 Companies | 50+ KPIs\n\n"
                     "Prepared by: Hashmil Muhammed\n"
                     "AI Engineer & Data Analyst | SCMS (SSET)\n"
                     "hashmilmuhammedparammal@gmail.com")

    # Slide 2: Problem Statement
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    slide.shapes.title.text = "Problem Statement"
    tf = slide.placeholders[1].text_frame
    tf.text = "Why institutional-grade analytics is needed:"
    p = tf.add_paragraph(); p.text = "Fragmented Data: Financial data is scattered across multiple annual reports and statements."; p.level = 1
    p = tf.add_paragraph(); p.text = "Lack of Automation: Manual computation of complex ratios and CAGR is error-prone."; p.level = 1
    p = tf.add_paragraph(); p.text = "No Peer Benchmarking: Difficult to compare a company's performance against sector medians instantly."; p.level = 1
    p = tf.add_paragraph(); p.text = "Missing Sentiment Analysis: Qualitative factors from business descriptions are often ignored."; p.level = 1

    # Slide 3: Project Objectives
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    slide.shapes.title.text = "Project Objectives"
    tf = slide.placeholders[1].text_frame
    tf.text = "Key Goals of FinSight N100:"
    p = tf.add_paragraph(); p.text = "Build a robust ETL pipeline and Star Schema SQLite database."; p.level = 1
    p = tf.add_paragraph(); p.text = "Compute 50+ KPIs including ROCE, FCF, and Turnaround CAGR."; p.level = 1
    p = tf.add_paragraph(); p.text = "Implement ML K-Means clustering & NLP Sentiment Scoring."; p.level = 1
    p = tf.add_paragraph(); p.text = "Deliver a 9-screen interactive Streamlit Dashboard and FastAPI backend."; p.level = 1

    # Slide 4: Data Sourcing & Ingestion
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    slide.shapes.title.text = "Data Sourcing & Ingestion"
    tf = slide.placeholders[1].text_frame
    tf.text = "12 Enterprise Datasets Utilized:"
    p = tf.add_paragraph(); p.text = "7 Core Datasets: 10-Year P&L, Balance Sheet, Cash Flow, and Business Descriptions."; p.level = 1
    p = tf.add_paragraph(); p.text = "5 Supplementary Datasets: Market Cap, Peer Groups, Sectors, Stock Prices."; p.level = 1
    p = tf.add_paragraph(); p.text = "Scale: 92 NSE Nifty 100 Companies generating over 11,000+ historical data points."; p.level = 1

    # Slide 5: System Architecture & ETL
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    slide.shapes.title.text = "System Architecture & ETL Pipeline"
    tf = slide.placeholders[1].text_frame
    tf.text = "End-to-End Workflow:"
    p = tf.add_paragraph(); p.text = "Extract: Automated pandas ingestion from raw Excel files."; p.level = 1
    p = tf.add_paragraph(); p.text = "Transform: 16 Data Quality (DQ) Rules applied (PK/FK uniqueness, tally verifications)."; p.level = 1
    p = tf.add_paragraph(); p.text = "Load: Indexed Star Schema SQLite Database (nifty100.db)."; p.level = 1
    p = tf.add_paragraph(); p.text = "Serve: 16 modular FastAPI endpoints serving data to Streamlit in <0.5s."; p.level = 1

    # Slide 6: Tech Stack
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    slide.shapes.title.text = "Technology Stack"
    tf = slide.placeholders[1].text_frame
    tf.text = "Core Technologies Used:"
    p = tf.add_paragraph(); p.text = "Backend & API: Python 3.10, FastAPI, Uvicorn"; p.level = 1
    p = tf.add_paragraph(); p.text = "Database: SQLite3 (Optimized & Indexed)"; p.level = 1
    p = tf.add_paragraph(); p.text = "Frontend: Streamlit, Plotly, Matplotlib"; p.level = 1
    p = tf.add_paragraph(); p.text = "Data Engineering & ML: Pandas, NumPy, Scikit-Learn (K-Means)"; p.level = 1
    p = tf.add_paragraph(); p.text = "NLP & QA: NLTK VADER, Pytest (137/137 passed)"; p.level = 1

    # Slide 7: Advanced Financial Analytics
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    slide.shapes.title.text = "Advanced Financial Analytics"
    tf = slide.placeholders[1].text_frame
    tf.text = "KPI Engine Highlights:"
    p = tf.add_paragraph(); p.text = "Negative-Base CAGR Resilience: Custom logic to handle turnaround financial years without math errors."; p.level = 1
    p = tf.add_paragraph(); p.text = "Banking Carve-outs: Adjusted ROCE and D/E logic specifically for financial sector constituents."; p.level = 1
    p = tf.add_paragraph(); p.text = "Cash Flow Intelligence: Capital Allocation mapping (8 phases like Cash Cow, Aggressive Growth)."; p.level = 1

    # Slide 8: Machine Learning & NLP
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    slide.shapes.title.text = "Machine Learning & NLP"
    tf = slide.placeholders[1].text_frame
    tf.text = "Adding AI Intelligence:"
    p = tf.add_paragraph(); p.text = "Unsupervised ML: K-Means Clustering (k=5) based on ROE, D/E, OPM, and CAGR for automated peer profiling."; p.level = 1
    p = tf.add_paragraph(); p.text = "NLP Tagging: Keyword-matching algorithm to classify sectors from raw text."; p.level = 1
    p = tf.add_paragraph(); p.text = "Sentiment Scoring: NLTK VADER applied to validate qualitative Pros/Cons generated by the rule engine."; p.level = 1

    # Slide 9: Smart Screener & Benchmarking
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    slide.shapes.title.text = "Smart Screener & Benchmarking"
    tf = slide.placeholders[1].text_frame
    tf.text = "Evaluating Market Leaders:"
    p = tf.add_paragraph(); p.text = "YAML Screener: Multi-criteria filtering engine (Quality, Value, Growth)."; p.level = 1
    p = tf.add_paragraph(); p.text = "Composite Ranking: Weighted scoring (Profitability 50%, Growth 30%, Valuation 20%)."; p.level = 1
    p = tf.add_paragraph(); p.text = "Radar Charts: 56 automated Plotly radar charts for sector-relative peer comparisons."; p.level = 1

    # Slide 10: Interactive Dashboard
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    slide.shapes.title.text = "Interactive Dashboard (Streamlit)"
    tf = slide.placeholders[1].text_frame
    tf.text = "9-Screen Business Intelligence UI:"
    p = tf.add_paragraph(); p.text = "Home / Overview: Market health and sector distribution."; p.level = 1
    p = tf.add_paragraph(); p.text = "Company Profile: Ticker search with 6-metric KPI tiles."; p.level = 1
    p = tf.add_paragraph(); p.text = "Trend Analysis: 10-year sparkline trajectories."; p.level = 1
    p = tf.add_paragraph(); p.text = "Automated Reports: PDF Tearsheet generation via ReportLab."; p.level = 1

    # Slide 11: Strategic Recommendations & Future Scope
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    slide.shapes.title.text = "Future Roadmap"
    tf = slide.placeholders[1].text_frame
    tf.text = "Scaling the Platform:"
    p = tf.add_paragraph(); p.text = "Docker Containerization: Package FastAPI and Streamlit for cloud deployment."; p.level = 1
    p = tf.add_paragraph(); p.text = "Real-Time Integration: WebSocket integration for live stock price streaming."; p.level = 1
    p = tf.add_paragraph(); p.text = "Advanced NLP: Upgrade to FinBERT for deep analysis of annual reports and earnings calls."; p.level = 1
    p = tf.add_paragraph(); p.text = "Alert System: Automated email triggers for screened companies."; p.level = 1

    # Slide 12: Thank You
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    title = slide.shapes.title
    subtitle = slide.placeholders[1]
    title.text = "Thank You"
    subtitle.text = ("Hashmil Muhammed\n"
                     "AI Engineer & Data Analyst\n"
                     "MCA Scholar, SCMS School of Engineering and Technology (SSET)\n"
                     "LinkedIn: linkedin.com/in/hashmil-muhammed08\n"
                     "GitHub: github.com/Hashmil-Muhammed")

    prs.save('FinSight_N100_Presentation.pptx')
    print("Presentation generated successfully: FinSight_N100_Presentation.pptx")

if __name__ == '__main__':
    create_presentation()